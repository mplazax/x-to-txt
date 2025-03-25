#!/usr/bin/env python3
import argparse
import os
import sys
from typing import List, Optional, Tuple

# Import required libraries for text extraction
try:
    from pdfminer.high_level import extract_text as extract_pdf_text
except ImportError:
    def extract_pdf_text(file_path):
        print("Error: pdfminer.six not installed. Run 'pip install pdfminer.six'.")
        sys.exit(1)

try:
    import docx
except ImportError:
    def open_docx(file_path):
        print("Error: python-docx not installed. Run 'pip install python-docx'.")
        sys.exit(1)
else:
    def open_docx(file_path):
        return docx.Document(file_path)

try:
    from pptx import Presentation
except ImportError:
    def open_pptx(file_path):
        print("Error: python-pptx not installed. Run 'pip install python-pptx'.")
        sys.exit(1)
else:
    def open_pptx(file_path):
        return Presentation(file_path)


def extract_text_from_pdf(file_path: str, verbose: bool = False) -> str:
    """Extract text from a PDF file."""
    try:
        text = extract_pdf_text(file_path)
        if verbose:
            # Count pages (approximation)
            with open(file_path, 'rb') as f:
                page_count = 0
                for line in f:
                    if b'/Page ' in line:
                        page_count += 1
            print(f"Processed {page_count} pages from PDF file: {file_path}")
        return text
    except Exception as e:
        print(f"Error extracting text from PDF: {str(e)}")
        return ""


def extract_text_from_docx(file_path: str, verbose: bool = False) -> str:
    """Extract text from a DOCX file."""
    try:
        doc = open_docx(file_path)
        paragraphs = [p.text for p in doc.paragraphs]
        text = "\n".join(paragraphs)
        
        if verbose:
            print(f"Processed {len(paragraphs)} paragraphs from DOCX file: {file_path}")
        
        return text
    except Exception as e:
        print(f"Error extracting text from DOCX: {str(e)}")
        return ""


def extract_text_from_pptx(file_path: str, verbose: bool = False) -> str:
    """Extract text from a PPTX file."""
    try:
        presentation = open_pptx(file_path)
        text_parts = []
        
        slide_count = len(presentation.slides)
        text_shape_count = 0
        
        for i, slide in enumerate(presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                    text_shape_count += 1
            
            if slide_text:
                text_parts.append(f"--- Slide {i} ---")
                text_parts.extend(slide_text)
        
        if verbose:
            print(f"Processed {slide_count} slides and {text_shape_count} text elements from PPTX file: {file_path}")
        
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting text from PPTX: {str(e)}")
        return ""


def determine_file_format(file_path: str, specified_format: Optional[str] = None) -> Optional[str]:
    """Determine the format of the file based on extension or specified format."""
    if specified_format:
        specified_format = specified_format.lower()
        if specified_format in ["pdf", "docx", "pptx"]:
            return specified_format
        else:
            print(f"Error: Invalid format specified: {specified_format}")
            return None
    
    _, extension = os.path.splitext(file_path)
    extension = extension.lower()
    
    if extension == ".pdf":
        return "pdf"
    elif extension == ".docx":
        return "docx"
    elif extension == ".pptx":
        return "pptx"
    else:
        print(f"Error: Unsupported file format for {file_path}")
        return None


def extract_text(file_path: str, format_type: Optional[str] = None, verbose: bool = False) -> Optional[str]:
    """Extract text from a file based on its format."""
    if not os.path.exists(file_path):
        print(f"Error: File not found: {file_path}")
        return None
    
    file_format = determine_file_format(file_path, format_type)
    if not file_format:
        return None
    
    if file_format == "pdf":
        return extract_text_from_pdf(file_path, verbose)
    elif file_format == "docx":
        return extract_text_from_docx(file_path, verbose)
    elif file_format == "pptx":
        return extract_text_from_pptx(file_path, verbose)
    else:
        # This should not happen given the checks in determine_file_format
        print(f"Error: Unsupported file format for {file_path}")
        return None


def handle_output(text: str, output_path: Optional[str], input_path: str) -> None:
    """Handle the output of extracted text."""
    if not text:
        return
    
    if not output_path:
        # Print to console
        print(text)
        return
    
    # Determine if output_path is a directory
    if os.path.isdir(output_path) or (not os.path.exists(output_path) and output_path.endswith(os.sep)):
        # Create directory if it doesn't exist
        os.makedirs(output_path, exist_ok=True)
        
        # Generate output file name based on input file name
        base_name = os.path.basename(input_path)
        name_without_ext = os.path.splitext(base_name)[0]
        output_file = os.path.join(output_path, f"{name_without_ext}.txt")
    else:
        # Ensure the directory for the output file exists
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
        
        output_file = output_path
    
    # Write text to the output file
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"Text extracted to: {output_file}")
    except Exception as e:
        print(f"Error writing to output file: {str(e)}")


def main() -> None:
    """Main function to handle command-line arguments and process files."""
    parser = argparse.ArgumentParser(
        description="Extract text from PDF, DOCX, and PPTX files.",
        epilog="Example: extract_text -o output_dir presentation.pptx document.pdf",
    )
    
    parser.add_argument("file_paths", nargs="+", help="One or more paths to the files from which to extract text")
    parser.add_argument(
        "-f", "--format", 
        choices=["pdf", "docx", "pptx"], 
        help="Specify the file format (pdf, docx, pptx). If omitted, format is inferred from file extension."
    )
    parser.add_argument(
        "-o", "--output", 
        help="Specify an output file (for a single input) or directory (for multiple inputs). If not provided, text is printed to the console."
    )
    parser.add_argument(
        "-v", "--verbose", 
        action="store_true", 
        help="Enable verbose output, displaying additional details like the number of pages or slides processed."
    )
    
    args = parser.parse_args()
    
    # Handle single file vs multiple files for output
    is_single_file = len(args.file_paths) == 1
    output_is_dir = args.output and (
        os.path.isdir(args.output) or 
        args.output.endswith(os.sep) or 
        not is_single_file
    )
    
    # Process each file
    for file_path in args.file_paths:
        if args.verbose:
            print(f"Processing file: {file_path}")
        
        text = extract_text(file_path, args.format, args.verbose)
        
        if text is not None:
            # Determine output path
            output_path = None
            if args.output:
                if output_is_dir:
                    # Create directory if needed
                    os.makedirs(args.output, exist_ok=True)
                    base_name = os.path.basename(file_path)
                    name_without_ext = os.path.splitext(base_name)[0]
                    output_path = os.path.join(args.output, f"{name_without_ext}.txt")
                else:
                    output_path = args.output
            
            handle_output(text, output_path, file_path)


if __name__ == "__main__":
    main() 